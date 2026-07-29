"""Slide decks: one page per slide, text only.

The rules, pinned: one `ExtractedPage` per slide, in deck order, so the
locator `pN` is slide N. A page is titled by the slide's title placeholder,
truncated to 80 characters, else "slide N". Within a slide, shapes are read
in shape XML order (groups recursed in place): each text frame with text
becomes one block whose lines are its paragraphs, and each table renders as
a markdown pipe table — first row as the header — in place. After the body,
a slide with speaker notes gets one final block prefixed exactly "Notes: ".
Blocks are separated by blank lines, the boundary the chunker splits on.

What this deliberately does not capture: charts, pictures, SmartArt, and
everything else that makes a deck a deck. The CLI says so at ingest and
points at the export-to-PDF + vision-extractor path, which does capture
them; text-only is the keyless floor, not the recommended path for a
visual-heavy deck.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import ExtractedPage, Extractor, ExtractionError, register

__all__ = ["PptxExtractor", "EXTRACTOR"]

_MEDIA_TYPES = frozenset({"pptx"})
_SUFFIXES = frozenset({".pptx"})

MAX_TITLE_CHARS = 80
"""A page title is its slide title, truncated to this."""


class PptxExtractor:
    """python-pptx reader. Deterministic: shape text in, slide pages out."""

    name = "pptx"
    version = "1"
    deterministic = True

    def can_handle(self, path: Path, media_type: str) -> bool:
        return media_type in _MEDIA_TYPES or path.suffix.lower() in _SUFFIXES

    def extract(self, path: Path, config: dict) -> Iterator[ExtractedPage]:
        try:
            deck = Presentation(str(path))
        except Exception as error:  # noqa: BLE001 - python-pptx raises broadly
            raise ExtractionError(f"could not open {path}: {error}") from error
        for number, slide in enumerate(deck.slides, start=1):
            blocks = list(_blocks(slide.shapes))
            notes = _notes(slide)
            if notes:
                blocks.append(f"Notes: {notes}")
            yield ExtractedPage(
                number=number,
                kind="page",
                name=_title(slide, number),
                text="\n\n".join(blocks),
            )


def _title(slide: Any, number: int) -> str:
    title = slide.shapes.title
    text = " ".join(title.text.split()) if title is not None else ""
    return text[:MAX_TITLE_CHARS] if text else f"slide {number}"


def _blocks(shapes: Any) -> Iterator[str]:
    """Block text per shape, in shape XML order, recursing into groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _blocks(shape.shapes)
            continue
        if getattr(shape, "has_table", False):
            rendered = _table(shape.table)
            if rendered:
                yield rendered
            continue
        if shape.has_text_frame:
            lines = [
                " ".join(paragraph.text.split())
                for paragraph in shape.text_frame.paragraphs
            ]
            block = "\n".join(line for line in lines if line)
            if block:
                yield block


def _notes(slide: Any) -> str:
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    return frame.text.strip() if frame is not None else ""


def _table(table: Any) -> str:
    """A markdown pipe table, first row as the header."""
    rows = [
        [_cell(cell.text) for cell in row.cells]
        for row in table.rows
    ]
    if not rows:
        return ""
    lines = ["| " + " | ".join(rows[0]) + " |"]
    lines.append("|" + "|".join(["---"] * len(rows[0])) + "|")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _cell(text: str) -> str:
    """A table cell on one line, with markdown's `|` neutralized."""
    return " ".join(text.split()).replace("|", "\\|")


EXTRACTOR: Extractor = PptxExtractor()
register(EXTRACTOR)
