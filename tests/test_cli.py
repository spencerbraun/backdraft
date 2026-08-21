"""The CLI surface W1 owns: discovery, session resolution, init/ingest/ls/export."""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from typer.testing import CliRunner

from backdraft import cli
from backdraft.extract.base import ExtractedPage
from backdraft.registry import DIRECTORY, Registry

runner = CliRunner()


@pytest.fixture(autouse=True)
def no_home_override(monkeypatch: pytest.MonkeyPatch) -> None:
    """Neither env var leaks in from the developer's shell."""
    monkeypatch.delenv(cli.HOME_ENV, raising=False)
    monkeypatch.delenv(cli.SESSION_ENV, raising=False)


@pytest.fixture
def project(tmp_path: Path, note: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """An initialized project, with cwd inside it."""
    monkeypatch.chdir(tmp_path)
    runner.invoke(cli.app, ["init"])
    return tmp_path


# ---- discovery --------------------------------------------------------------


def test_find_root_walks_up(tmp_path: Path) -> None:
    (tmp_path / DIRECTORY).mkdir()
    deep = tmp_path / "a" / "b"
    deep.mkdir(parents=True)
    assert cli.find_root(deep) == tmp_path.resolve()


def test_find_root_is_none_without_a_registry(tmp_path: Path) -> None:
    assert cli.find_root(tmp_path) is None


def test_backdraft_home_overrides_the_walk(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setenv(cli.HOME_ENV, str(tmp_path))
    assert cli.find_root(Path("/")) == tmp_path


def test_backdraft_home_accepts_the_registry_directory_itself(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setenv(cli.HOME_ENV, str(tmp_path / DIRECTORY))
    assert cli.find_root(Path("/")) == tmp_path


# ---- sessions ---------------------------------------------------------------


def test_the_default_session_is_used_when_nothing_is_given(root: Path) -> None:
    with Registry.open(root) as registry:
        assert cli.resolve_session(registry=registry) == "default"


def test_the_env_var_beats_the_default(root: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv(cli.SESSION_ENV, "from-env")
    with Registry.open(root) as registry:
        assert cli.resolve_session(registry=registry) == "from-env"


def test_the_flag_beats_the_env_var(root: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv(cli.SESSION_ENV, "from-env")
    with Registry.open(root) as registry:
        assert cli.resolve_session("from-flag", registry) == "from-flag"


# ---- commands ---------------------------------------------------------------


def test_init_creates_the_registry(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.chdir(tmp_path)
    result = runner.invoke(cli.app, ["init"])
    assert result.exit_code == 0
    assert (tmp_path / DIRECTORY / "registry.db").is_file()
    assert "documents: 0" in result.stdout


def test_commands_without_a_registry_exit_one(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.chdir(tmp_path)
    result = runner.invoke(cli.app, ["ls"])
    assert result.exit_code == cli.EXIT_USAGE
    assert "backdraft init" in result.stderr


def test_ingest_then_ls(project: Path, note: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(note)])
    assert result.exit_code == 0
    assert "quarterly-notes" in result.stdout

    listed = runner.invoke(cli.app, ["ls"])
    assert listed.exit_code == 0
    assert "quarterly-notes\tquarterly-notes.md\ttext\t1 page" in listed.stdout


def test_ls_says_so_when_empty(project: Path) -> None:
    assert "no documents ingested" in runner.invoke(cli.app, ["ls"]).stdout


def test_ls_and_read_count_a_workbook_in_the_same_words(
    project: Path, workbook: Path
) -> None:
    """One registry, one vocabulary.

    `ls` said `2 pages` where the gate's list said `2 sheets`, which reads as
    two commands describing two different registries. Both now go through
    `reader.unit`; this pins them together rather than pinning either string.
    """
    runner.invoke(cli.app, ["ingest", str(workbook)])
    listed = runner.invoke(cli.app, ["ls"]).stdout
    read = runner.invoke(cli.app, ["read"]).stdout
    assert "2 sheets" in listed and "2 sheets" in read
    assert "2 pages" not in listed


def test_a_single_page_source_is_counted_in_the_singular(
    project: Path, note: Path
) -> None:
    """`1 pages` is the common case for one-page sources, and it was ungrammatical."""
    ingested = runner.invoke(cli.app, ["ingest", str(note)]).stdout
    assert "1 page" in ingested and "1 pages" not in ingested
    assert "1 pages" not in runner.invoke(cli.app, ["ls"]).stdout


def test_ingest_accepts_an_explicit_slug(project: Path, note: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(note), "--slug", "t12-audit"])
    assert result.exit_code == 0
    assert "t12-audit" in result.stdout


def test_a_slug_with_several_files_is_a_usage_error(project: Path, note: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(note), str(note), "--slug", "x"])
    assert result.exit_code == cli.EXIT_USAGE


def test_ingest_reports_an_extractor_failure(project: Path, tmp_path: Path) -> None:
    broken = tmp_path / "broken.pdf"
    broken.write_bytes(b"%PDF-1.4 nonsense")
    result = runner.invoke(cli.app, ["ingest", str(broken)])
    assert result.exit_code == cli.EXIT_USAGE
    assert "broken.pdf" in result.stderr


# ---- ingest finishes its list -----------------------------------------------


# Long enough to be a real source: under THIN_SOURCE_CHARS these would each
# carry the thin-source note, which is true of a one-line file and beside the
# point of the tests below.
_A = """\
Occupancy closed at 91.4% in Q3, up sixty basis points on the quarter and the
third consecutive quarter of gains. Concessions burned off across the Riverside
and Kenwood assets, and the lease-up at Riverside reached stabilization in
August, a month ahead of the sponsor's underwriting.
"""
_C = """\
Debt service coverage was 1.28x for the trailing twelve months, clearing the
1.25x covenant by a margin the lender described as adequate but not comfortable.
The March refinancing lowered the blended coupon by roughly seventy basis points,
which is most of what separates this quarter's coverage from last year's.
"""


def _sources(tmp_path: Path) -> tuple[Path, Path, Path]:
    """Two readable markdown files with an unreadable PDF between them."""
    first = tmp_path / "a.md"
    first.write_text(_A, encoding="utf-8")
    broken = tmp_path / "broken.pdf"
    broken.write_bytes(b"%PDF-1.4 nonsense")
    last = tmp_path / "c.md"
    last.write_text(_C, encoding="utf-8")
    return first, broken, last


def test_a_failed_source_does_not_stop_the_ones_behind_it(
    project: Path, tmp_path: Path
) -> None:
    """The case this exists for: `c.md` used never to be attempted at all."""
    first, broken, last = _sources(tmp_path)
    result = runner.invoke(cli.app, ["ingest", str(first), str(broken), str(last)])
    assert result.exit_code == cli.EXIT_USAGE
    assert "a-doc" in result.stdout and "c-doc" in result.stdout
    failures = [line for line in result.stderr.splitlines() if line.startswith("  ! ")]
    assert len(failures) == 1
    assert "broken.pdf" in failures[0] and "PDF" in failures[0]
    assert "2 of 3 sources ingested; 1 failed" in result.stderr
    listed = runner.invoke(cli.app, ["ls"]).stdout
    assert "a-doc" in listed and "c-doc" in listed


def test_the_failure_report_says_re_running_the_list_is_safe(
    project: Path, tmp_path: Path
) -> None:
    """Without this the caller has to diff `ls` against its own arguments to
    work out which half to pass again."""
    first, broken, last = _sources(tmp_path)
    result = runner.invoke(cli.app, ["ingest", str(first), str(broken), str(last)])
    assert "re-run the same command" in result.stderr
    assert "no-op" in result.stderr


def test_every_source_failing_still_exits_one_and_names_each(
    project: Path, tmp_path: Path
) -> None:
    broken = tmp_path / "broken.pdf"
    broken.write_bytes(b"%PDF-1.4 nonsense")
    also = tmp_path / "also-broken.pdf"
    also.write_bytes(b"%PDF-1.4 nonsense either way")
    result = runner.invoke(cli.app, ["ingest", str(broken), str(also)])
    assert result.exit_code == cli.EXIT_USAGE
    assert "0 of 2 sources ingested; 2 failed" in result.stderr
    assert "broken.pdf" in result.stderr and "also-broken.pdf" in result.stderr
    with Registry.open(project) as registry:
        assert registry.documents() == []


def test_one_source_failing_alone_counts_in_the_singular(
    project: Path, tmp_path: Path
) -> None:
    broken = tmp_path / "broken.pdf"
    broken.write_bytes(b"%PDF-1.4 nonsense")
    result = runner.invoke(cli.app, ["ingest", str(broken)])
    assert "0 of 1 source ingested; 1 failed" in result.stderr
    assert "1 sources" not in result.stderr


def test_a_run_where_nothing_fails_prints_one_line_per_source_and_nothing_else(
    project: Path, tmp_path: Path
) -> None:
    """The failure report is failure-only, and so is every note.

    A clean run of readable, ordinary sources gains no prefix, no trailing note
    and no second line per source — each source's whole story is its own line.
    Pinned as an exact string because every note added here is a note that could
    have leaked into it."""
    first, _, last = _sources(tmp_path)
    result = runner.invoke(cli.app, ["ingest", str(first), str(last)])
    assert result.exit_code == 0
    assert result.stdout == (
        "a-doc  a.md  text  1 page  286 chars\n"
        "c-doc  c.md  text  1 page  314 chars\n"
    )
    assert result.stderr == ""


# ---- ingest says what it did and what it got --------------------------------


def test_a_fresh_document_reports_how_much_text_came_out(
    project: Path, note: Path
) -> None:
    """The count that says whether the snapshot is worth citing at all."""
    result = runner.invoke(cli.app, ["ingest", str(note)])
    assert result.exit_code == 0
    line = result.stdout.splitlines()[0]
    assert line.startswith("quarterly-notes  quarterly-notes.md  text  1 page  ")
    assert line.endswith(" chars")
    chars = int(line.rsplit("  ", 1)[1].removesuffix(" chars"))
    assert chars == len(note.read_text(encoding="utf-8"))


def test_re_ingesting_unchanged_bytes_says_so(project: Path, note: Path) -> None:
    """A no-op printed like a fresh ingest is how an agent loses track of state."""
    first = runner.invoke(cli.app, ["ingest", str(note)])
    assert "unchanged" not in first.stdout

    second = runner.invoke(cli.app, ["ingest", str(note)])
    assert second.exit_code == 0
    assert second.stdout.splitlines()[0].endswith("  unchanged")
    assert "new generation" not in second.stdout


def test_re_ingesting_edited_bytes_says_a_generation_was_made_and_names_drift(
    project: Path, note: Path
) -> None:
    """The moment older citations can start reporting `drifted` — said out loud."""
    runner.invoke(cli.app, ["ingest", str(note)])
    note.write_text(
        note.read_text(encoding="utf-8").replace("1.42x", "1.51x"), encoding="utf-8"
    )
    result = runner.invoke(cli.app, ["ingest", str(note)])
    assert result.exit_code == 0
    assert result.stdout.splitlines()[0].endswith("  new generation")
    assert "note: new generation of quarterly-notes" in result.stdout
    assert "`drifted`" in result.stdout
    assert "backdraft bind" in result.stdout


def test_one_note_names_every_document_that_gained_a_generation(
    project: Path, note: Path, workbook: Path
) -> None:
    """Grouped, not one line per document: the consequence is the same for each."""
    runner.invoke(cli.app, ["ingest", str(note), str(workbook)])
    note.write_text(
        note.read_text(encoding="utf-8").replace("1.42x", "1.51x"), encoding="utf-8"
    )
    from openpyxl import load_workbook

    book = load_workbook(workbook)
    book.active["A1"] = "Unit no."
    book.save(workbook)

    result = runner.invoke(cli.app, ["ingest", str(note), str(workbook)])
    assert result.exit_code == 0
    notes = [line for line in result.stdout.splitlines() if "new generation of" in line]
    assert len(notes) == 1
    assert "quarterly-notes, model" in notes[0]


def test_a_thin_source_says_the_likely_cause_and_what_to_do(
    project: Path, tmp_path: Path
) -> None:
    """A login wall ingests cleanly, exits 0, and used to print `1 page` like any
    success — so an agent could cite the shell of a source without a signal."""
    wall = tmp_path / "q4-results.html"
    wall.write_text(
        "<html><body><p>Please sign in to continue.</p></body></html>",
        encoding="utf-8",
    )
    result = runner.invoke(cli.app, ["ingest", str(wall)])
    assert result.exit_code == 0, "a thin source is a real snapshot, never a failure"
    assert "note: little text extracted" in result.stdout
    assert "login wall" in result.stdout
    assert "backdraft read" in result.stdout
    assert "q4-results" in result.stdout.splitlines()[-1]


def test_a_thin_pdf_is_told_about_the_missing_text_layer_instead(
    project: Path, tmp_path: Path, scripted: type
) -> None:
    """The cause a media type can actually say: a scan has no text layer to read."""
    scan = tmp_path / "scan.pdf"
    scan.write_bytes(b"%PDF-1.4 stub")
    scripted("thinpdf", [ExtractedPage(number=1, kind="page", text="Exhibit A")])
    result = runner.invoke(cli.app, ["ingest", str(scan), "--extractor", "thinpdf"])
    assert result.exit_code == 0
    assert "note: little text extracted" in result.stdout
    assert "no text layer is a scan" in result.stdout
    assert "BACKDRAFT_VLM_API_KEY" in result.stdout
    assert "login wall" not in result.stdout


def test_a_source_that_extracted_nothing_at_all_is_thin_too(
    project: Path, tmp_path: Path, scripted: type
) -> None:
    """Zero pages is the limit case of thin, and the one a scan actually hits."""
    empty = tmp_path / "blank.pdf"
    empty.write_bytes(b"%PDF-1.4 stub")
    scripted("nopages", [])
    result = runner.invoke(cli.app, ["ingest", str(empty), "--extractor", "nopages"])
    assert result.exit_code == 0
    assert result.stdout.splitlines()[0] == "blank  blank.pdf  pdf  0 pages  0 chars"
    assert "note: little text extracted" in result.stdout


def test_a_normal_source_carries_no_thin_note(project: Path, note: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(note)])
    assert result.exit_code == 0
    assert "little text extracted" not in result.stdout


def test_a_thin_deck_is_pointed_at_the_note_that_already_named_the_gap(
    project: Path, tmp_path: Path
) -> None:
    """A deck of images is the one thin cause the command already prints.

    Falling through to the generic default would tell an agent less than the
    line directly above it does, so `pptx` names the cause and defers the fix
    rather than restating it.
    """
    from pptx import Presentation

    deck = tmp_path / "visual.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(deck)

    result = runner.invoke(cli.app, ["ingest", str(deck)])
    assert result.exit_code == 0
    assert "note: extracted slide text only" in result.stdout
    assert "note: little text extracted — a deck whose slides are" in result.stdout
    assert "the note above has the fix" in result.stdout
    assert "may simply be short" not in result.stdout


def test_a_thin_source_with_nothing_specific_to_say_says_so_plainly(
    project: Path, tmp_path: Path
) -> None:
    """A media type off the cause table gets the default, not a wrong guess.

    `text` has no story a scan or a login wall has: an almost-empty note is
    almost certainly an almost-empty note. The note still fires, because the
    thing worth saying — do not cite the shell of this — is the same.
    """
    stub = tmp_path / "stub.md"
    stub.write_text("Occupancy was fine.\n", encoding="utf-8")
    result = runner.invoke(cli.app, ["ingest", str(stub)])
    assert result.exit_code == 0
    assert "note: little text extracted — the source may simply be short" in result.stdout
    assert "no text layer" not in result.stdout
    assert "login wall" not in result.stdout
    assert result.stdout.rstrip().endswith("stub.")


def test_ingesting_a_deck_notes_the_text_only_gap(project: Path, tmp_path: Path) -> None:
    """The note a calling agent relays when the deck is visual-heavy."""
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Q3"
    path = tmp_path / "q3-deck.pptx"
    deck.save(str(path))
    result = runner.invoke(cli.app, ["ingest", str(path)])
    assert result.exit_code == 0
    assert "q3-deck" in result.stdout
    assert (
        "note: extracted slide text only. Charts and images on slides are "
        "not captured; exporting the deck to PDF and ingesting it through "
        "the vision extractor captures them." in result.stdout
    )


def test_ingesting_a_text_file_carries_no_slide_note(project: Path, note: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(note)])
    assert "slide text" not in result.stdout


def test_ingest_rejects_a_malformed_config_pair(project: Path, note: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(note), "--config", "nonsense"])
    assert result.exit_code == cli.EXIT_USAGE
    assert "key=value" in result.stderr


def test_ingest_rejects_a_key_the_chosen_extractor_never_reads(
    project: Path, note: Path
) -> None:
    """A markdown note goes through `text`, which reads nothing — so the typo
    that motivated this and the key that is real elsewhere fail the same way."""
    result = runner.invoke(cli.app, ["ingest", str(note), "--config", "dpi=300"])
    assert result.exit_code == cli.EXIT_USAGE
    assert "unknown config key 'dpi' for text" in result.stderr
    assert "reads no config keys" in result.stderr
    with Registry.open(project) as registry:
        assert registry.documents() == []


def test_a_rejected_config_key_names_the_ones_that_would_have_worked(
    project: Path, note: Path, scripted: type
) -> None:
    scripted(
        "tunable",
        [ExtractedPage(number=1, kind="page", text="one page, however it was tuned")],
        config_keys={"mode": "how loudly to read", "voice": "whose voice to read in"},
    )
    result = runner.invoke(
        cli.app, ["ingest", str(note), "--extractor", "tunable", "--config", "mdoe=loud"]
    )
    assert result.exit_code == cli.EXIT_USAGE
    assert "unknown config key 'mdoe' for tunable; known: mode, voice" in result.stderr


def test_ingest_passes_config_through(project: Path, note: Path, scripted: type) -> None:
    scripted(
        "tunable",
        [ExtractedPage(number=1, kind="page", text="one page, however it was tuned")],
        config_keys={"mode": "how loudly to read"},
    )
    first = runner.invoke(cli.app, ["ingest", str(note), "--extractor", "tunable"])
    second = runner.invoke(
        cli.app, ["ingest", str(note), "--extractor", "tunable", "--config", "mode=loud"]
    )
    assert (first.exit_code, second.exit_code) == (0, 0)
    with Registry.open(project) as registry:
        generations = registry.export_json()["documents"][0]["extractions"]
    assert len(generations) == 2


def test_export_to_stdout(project: Path, note: Path) -> None:
    runner.invoke(cli.app, ["ingest", str(note)])
    result = runner.invoke(cli.app, ["export"])
    assert result.exit_code == 0
    assert json.loads(result.stdout)["$format"] == "backdraft/registry-v1"


def test_export_to_a_file(project: Path, note: Path, tmp_path: Path) -> None:
    runner.invoke(cli.app, ["ingest", str(note)])
    out = tmp_path / "registry.json"
    result = runner.invoke(cli.app, ["export", "--out", str(out)])
    assert result.exit_code == 0
    assert json.loads(out.read_text())["documents"][0]["slug"] == "quarterly-notes"


def test_ingesting_a_workbook_through_the_cli(project: Path, workbook: Path) -> None:
    result = runner.invoke(cli.app, ["ingest", str(workbook)])
    assert result.exit_code == 0
    assert "model" in result.stdout and "2 sheets" in result.stdout


def test_help_lists_the_commands_w1_owns() -> None:
    result = runner.invoke(cli.app, ["--help"])
    for command in ("init", "ingest", "ls", "export"):
        assert command in result.stdout


def test_ingest_help_names_every_extractor() -> None:
    """The help is where a caller looks before guessing a name and getting an
    error. A registered extractor missing from it is a doc that lies."""
    from backdraft.extract import base

    help_text = runner.invoke(cli.app, ["ingest", "--help"]).stdout
    # Typer wraps the help column, so match on the un-wrapped run of words.
    flattened = " ".join(help_text.split())
    missing = [name for name in base.names() if name not in flattened]
    assert missing == [], missing


def test_every_command_summary_is_one_line() -> None:
    """The command table is scannable only if each summary fits its row: typer
    prints the docstring's whole first paragraph, so a two-line first sentence
    silently doubles a row's height."""
    over: list[str] = []
    for command in cli.app.registered_commands:
        summary = (command.callback.__doc__ or "").strip().split("\n\n")[0]
        if len(" ".join(summary.split())) > 80:
            over.append(command.callback.__name__)
    assert over == [], over


def test_a_missing_sub_app_degrades_silently() -> None:
    """A workstream that has not landed yet must not break the whole CLI."""
    assert cli._mount("backdraft.nothing.here.cli") is False
    assert runner.invoke(cli.app, ["--help"]).exit_code == 0


def test_skill_install_copies_the_bundled_skill(tmp_path, monkeypatch) -> None:
    from typer.testing import CliRunner

    from backdraft.cli import app

    monkeypatch.setattr("pathlib.Path.home", lambda: tmp_path)
    result = CliRunner().invoke(app, ["skill", "install"])
    assert result.exit_code == 0, result.output
    installed = tmp_path / ".claude" / "skills" / "backdraft" / "SKILL.md"
    assert installed.is_file()
    assert "backdraft read" in installed.read_text(encoding="utf-8")

    result = CliRunner().invoke(app, ["skill", "install", "--all"])
    assert result.exit_code == 0, result.output
    assert (tmp_path / ".claude" / "skills" / "backdraft-backfill" / "SKILL.md").is_file()


def test_skill_install_default_agent_touches_only_claude(tmp_path, monkeypatch) -> None:
    from typer.testing import CliRunner

    from backdraft.cli import app

    monkeypatch.setattr("pathlib.Path.home", lambda: tmp_path)
    result = CliRunner().invoke(app, ["skill", "install"])
    assert result.exit_code == 0, result.output
    assert (tmp_path / ".claude" / "skills" / "backdraft" / "SKILL.md").is_file()
    assert not (tmp_path / ".agents").exists()


def test_skill_install_agent_codex_uses_the_standard_path(tmp_path, monkeypatch) -> None:
    from typer.testing import CliRunner

    from backdraft.cli import app

    monkeypatch.setattr("pathlib.Path.home", lambda: tmp_path)
    result = CliRunner().invoke(app, ["skill", "install", "--agent", "codex"])
    assert result.exit_code == 0, result.output
    assert (tmp_path / ".agents" / "skills" / "backdraft" / "SKILL.md").is_file()
    assert not (tmp_path / ".claude").exists()


def test_skill_install_agent_all_lands_in_both_layouts(tmp_path, monkeypatch) -> None:
    from typer.testing import CliRunner

    from backdraft.cli import app

    monkeypatch.setattr("pathlib.Path.home", lambda: tmp_path)
    result = CliRunner().invoke(app, ["skill", "install", "--agent", "all", "--all"])
    assert result.exit_code == 0, result.output
    for root in (".claude", ".agents"):
        for name in ("backdraft", "backdraft-backfill", "backdraft-artifact"):
            assert (tmp_path / root / "skills" / name / "SKILL.md").is_file()


def test_skill_install_agent_codex_project_lands_in_cwd(tmp_path, monkeypatch) -> None:
    from typer.testing import CliRunner

    from backdraft.cli import app

    monkeypatch.chdir(tmp_path)
    result = CliRunner().invoke(app, ["skill", "install", "--agent", "codex", "--project"])
    assert result.exit_code == 0, result.output
    assert (tmp_path / ".agents" / "skills" / "backdraft" / "SKILL.md").is_file()


def test_skill_install_rejects_an_unknown_agent(tmp_path, monkeypatch) -> None:
    from typer.testing import CliRunner

    from backdraft.cli import app

    monkeypatch.setattr("pathlib.Path.home", lambda: tmp_path)
    result = CliRunner().invoke(app, ["skill", "install", "--agent", "cursor"])
    assert result.exit_code == 1
    assert "unknown agent" in result.output


def test_python_dash_m_backdraft_runs_the_cli() -> None:
    import subprocess
    import sys

    result = subprocess.run(
        [sys.executable, "-m", "backdraft", "--help"],
        capture_output=True,
        text=True,
    )
    assert result.returncode == 0, result.stderr
    assert "ingest" in result.stdout


def test_skill_descriptions_fit_the_upload_cap() -> None:
    """claude.ai's skill upload rejects descriptions over 200 characters."""
    import re

    from backdraft.cli import SKILLS, _skills_source

    for name in SKILLS:
        text = (_skills_source() / name / "SKILL.md").read_text(encoding="utf-8")
        match = re.search(r"^description: (.+)$", text, flags=re.MULTILINE)
        assert match, name
        assert len(match.group(1)) <= 200, (name, len(match.group(1)))
        # An unquoted YAML scalar dies on ": " — keep descriptions parseable.
        assert ": " not in match.group(1), name
        named = re.search(r"^name: (.+)$", text, flags=re.MULTILINE)
        assert named and named.group(1) == name
